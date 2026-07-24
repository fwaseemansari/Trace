from pptx import Presentation
from langchain_community.document_loaders import TextLoader, CSVLoader
from langchain_core.documents import Document
import pdfplumber
import os
import openpyxl
import re

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None


import re

def load_pdf(file_path):
    docs = []
    current_heading = None

    # Table-header-like phrases to exclude from heading detection —
    # these are column labels, not real section titles, even though
    # they're often styled bold/larger like headings.
    TABLE_HEADER_WORDS = {"day", "task", "owner", "specific", "output", "file", "item", "field", "description"}

    def looks_like_table_header(text):
        words = set(re.findall(r"[a-zA-Z]+", text.lower()))
        # If most of the words in the candidate heading are generic
        # table-column terms, treat it as a table header, not a real heading.
        if not words:
            return False
        overlap = words & TABLE_HEADER_WORDS
        return len(overlap) >= 2 and len(overlap) / len(words) > 0.4

    with pdfplumber.open(file_path) as pdf:
        all_sizes = []
        try:
            for page in pdf.pages:
                for char in page.chars:
                    all_sizes.append(round(char.get("size", 0)))
        except Exception as e:
            print(f"[load_pdf] font-size scan failed: {e}")
        body_size = max(set(all_sizes), key=all_sizes.count) if all_sizes else 0

        for page_num, page in enumerate(pdf.pages, start=1):
            try:
                lines = page.extract_text_lines() if hasattr(page, "extract_text_lines") else []
            except Exception as e:
                print(f"[load_pdf] extract_text_lines failed on page {page_num}: {e}")
                lines = []

            segments = []
            buffer = []

            for line in lines:
                chars = line.get("chars", [])
                text = line.get("text", "").strip()
                if not text:
                    continue

                is_heading = False
                if chars:
                    avg_size = sum(c.get("size", 0) for c in chars) / len(chars)
                    bold = any("Bold" in (c.get("fontname", "") or "") for c in chars)
                    if (avg_size > body_size + 2 or (bold and avg_size >= body_size)) and len(text) < 120:
                        is_heading = True

                # Reject false-positive headings that are actually table column headers
                if is_heading and looks_like_table_header(text):
                    is_heading = False

                if is_heading:
                    if buffer:
                        segments.append(("\n".join(buffer), current_heading))
                        buffer = []
                    current_heading = text
                else:
                    buffer.append(text)

            if buffer:
                segments.append(("\n".join(buffer), current_heading))

            if not segments:
                try:
                    plain_text = page.extract_text() or ""
                except Exception as e:
                    print(f"[load_pdf] extract_text failed on page {page_num}: {e}")
                    plain_text = ""
                if plain_text.strip():
                    segments = [(plain_text, current_heading)]

            table_text = ""
            try:
                tables = page.extract_tables()
                for table in tables:
                    if not table:
                        continue
                    headers = table[0]
                    for row in table[1:]:
                        if not any(row):
                            continue
                        row_text = ", ".join(
                            f"{str(headers[i]).strip()}: {str(cell).strip()}"
                            for i, cell in enumerate(row)
                            if cell and i < len(headers) and headers[i]
                        )
                        if row_text:
                            table_text += row_text + "\n"
            except Exception as e:
                print(f"[load_pdf] table extraction failed on page {page_num}: {e}")

            if table_text.strip():
                segments.append((table_text, current_heading))

            for seg_text, seg_heading in segments:
                if not seg_text.strip():
                    continue
                content = f"[Section: {seg_heading}]\n{seg_text}" if seg_heading else seg_text
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source": file_path,
                        "page": page_num - 1,
                        "page_label": str(page_num)
                    }
                ))

    print(f"[load_pdf] {file_path}: extracted {len(docs)} segment-documents")
    return docs

def load_docx(file_path):
    if DocxDocument is None:
        raise ImportError("python-docx not installed. Run: pip install python-docx")

    doc = DocxDocument(file_path)
    full_text = ""

    for para in doc.paragraphs:
        if para.text.strip():
            full_text += para.text + "\n"

    for table in doc.tables:
        if not table.rows:
            continue
        headers = [cell.text.strip() for cell in table.rows[0].cells]
        for row in table.rows[1:]:
            cells = [cell.text.strip() for cell in row.cells]
            if not any(cells):
                continue
            row_text = ", ".join(
                f"{headers[i]}: {cells[i]}"
                for i in range(min(len(headers), len(cells)))
                if cells[i] and headers[i]
            )
            if row_text:
                full_text += row_text + "\n"

    return [Document(
        page_content=full_text,
        metadata={"source": file_path, "page": 0, "page_label": "1"}
    )]
    
def load_pptx(file_path):
    prs = Presentation(file_path)
    full_text = ""
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        full_text += f"\n[Slide {slide_num}]\n"
        
        for shape in slide.shapes:
            # regular text boxes
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        full_text += para.text + "\n"
            
            # tables in slides
            if shape.has_table:
                table = shape.table
                rows = list(table.rows)
                headers = [cell.text.strip() for cell in rows[0].cells]
                for row in rows[1:]:
                    cells = [cell.text.strip() for cell in row.cells]
                    if not any(cells):
                        continue
                    row_text = ", ".join(
                        f"{headers[i]}: {cells[i]}"
                        for i in range(min(len(headers), len(cells)))
                        if cells[i] and headers[i]
                    )
                    if row_text:
                        full_text += row_text + "\n"
    
    return [Document(
        page_content=full_text,
        metadata={"source": file_path, "page": 0, "page_label": "1"}
    )]

def load_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path, data_only=True)
    docs = []

    for sheet_num, sheet_name in enumerate(wb.sheetnames, start=1):
        sheet = wb[sheet_name]
        rows = list(sheet.iter_rows(values_only=True))

        if not rows:
            continue

        headers = [str(cell).strip() if cell is not None else "" for cell in rows[0]]
        full_text = f"[Sheet: {sheet_name}]\n"

        for row in rows[1:]:
            if not any(cell is not None for cell in row):
                continue
            row_text = ", ".join(
                f"{headers[i]}: {str(cell).strip()}"
                for i, cell in enumerate(row)
                if i < len(headers) and headers[i] and cell is not None
            )
            if row_text:
                full_text += row_text + "\n"

        if full_text.strip():
            docs.append(Document(
                page_content=full_text,
                metadata={
                    "source": file_path,
                    "page": sheet_num - 1,
                    "page_label": sheet_name
                }
            ))

    return docs

def load_document(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return load_pdf(file_path)
    elif ext == ".docx":
        return load_docx(file_path)
    elif ext in (".txt", ".md"):
        return TextLoader(file_path, encoding="utf-8").load()
    elif ext == ".csv":
        return CSVLoader(file_path).load()
    elif ext == ".xlsx":
        return load_xlsx(file_path)
    elif ext == ".pptx":
        return load_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

